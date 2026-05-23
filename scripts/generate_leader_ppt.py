#!/usr/bin/env python3
"""Generate NDR 2.0 leadership briefing PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Theme colors
COLOR_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)   # dark blue
COLOR_ACCENT = RGBColor(0x2E, 0x75, 0xB6)    # medium blue
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_HIGHLIGHT = RGBColor(0xC0, 0x00, 0x00)


def set_run(run, size=18, bold=False, color=None, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    set_run(r, 32, True, COLOR_PRIMARY)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.2))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        set_run(r2, 16, False, COLOR_MUTED)
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.4))
    fp = footer.text_frame.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = "NDR 2.0 可行性调研 · 第一阶段  |  2026年5月"
    set_run(fr, 10, False, COLOR_MUTED)


def add_section_slide(prs, title, bullets, footer_note=None, highlight_first=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tp = tbox.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    set_run(tr, 24, True, RGBColor(0xFF, 0xFF, 0xFF))

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(12)
        if isinstance(item, tuple):
            label, text = item
            r0 = p.add_run()
            r0.text = label + " "
            set_run(r0, 16, True, COLOR_ACCENT if not highlight_first else COLOR_HIGHLIGHT)
            r1 = p.add_run()
            r1.text = text
            set_run(r1, 16, i == 0 and highlight_first, COLOR_TEXT)
        else:
            r = p.add_run()
            r.text = "• " + item
            set_run(r, 16, False, COLOR_TEXT)

    if footer_note:
        fb = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(8.8), Inches(0.5))
        fp = fb.text_frame.paragraphs[0]
        fr = fp.add_run()
        fr.text = footer_note
        set_run(fr, 9, False, COLOR_MUTED)


def add_table_slide(prs, title, headers, rows, footer_note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tr = tbox.text_frame.paragraphs[0].add_run()
    tr.text = title
    set_run(tr, 24, True, RGBColor(0xFF, 0xFF, 0xFF))

    cols, row_count = len(headers), len(rows) + 1
    table = slide.shapes.add_table(row_count, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.35 * row_count)).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_ACCENT
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run(r, 11, True, RGBColor(0xFF, 0xFF, 0xFF))
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    set_run(r, 10, False, COLOR_TEXT)

    if footer_note:
        fb = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        fr = fb.text_frame.paragraphs[0].add_run()
        fr.text = footer_note
        set_run(fr, 9, False, COLOR_MUTED)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1 - Cover
    add_title_slide(
        prs,
        "NDR 2.0 加密流量分析",
        "需求可行性验证（第一阶段）\n汇报人：________    日期：2026年5月",
    )

    # Slide 2 - Conclusion
    add_section_slide(
        prs,
        "一句话结论",
        [
            ("建议决策：", "批准 NDR 2.0 加密流量分析纳入版本规划"),
            ("需求成立：", "87.2% 攻击经加密通道；IDC 将加密流量检测列为 2024 核心趋势"),
            ("市场可观：", "中国 NDR 23.6 亿元（2024）；全球 37.7 亿美元（2025）"),
            ("方向清晰：", "AI 不解密检测为行业共识，竞品已布局，存在差异化窗口"),
        ],
        footer_note="依据：NDR市场加密流量检测深度分析报告 V1.1（证据标注版）",
        highlight_first=True,
    )

    # Slide 3 - Why
    add_section_slide(
        prs,
        "为什么要做？— 威胁驱动 + 技术倒逼",
        [
            ("加密攻击占比 ", "87.2%（Zscaler 321 亿样本，同比 +10.3%）"),
            ("恶意软件占比 ", "86.5% 加密威胁为恶意软件（勒索/木马/后门）"),
            ("HTTPS 普及 ", "98% Web 请求已加密（HTTP Archive 2024）"),
            ("传统手段失效 ", "明文 DPI 无法有效检测 TLS 1.3 加密流量（73% 站点已采用）"),
        ],
        footer_note="📸 证据截图：Zscaler 博客 S-04 | HTTP Archive S-05",
    )

    # Slide 4 - Market size
    add_table_slide(
        prs,
        "市场有多大？— 赛道规模支撑投入",
        ["市场", "规模", "增速/预测", "来源"],
        [
            ["中国 NDR", "23.6 亿元人民币", "2024 年 -0.9%", "IDC CHC52196225"],
            ["全球 NDR", "37.7 亿美元", "2025 年基准", "Grand View Research"],
            ["全球 NDR", "58.2 亿美元", "2030 年，CAGR 9.6%", "MarketsandMarkets"],
            ["技术趋势", "加密流量检测", "IDC 2024 核心趋势", "IDC PR"],
        ],
        footer_note="📸 证据截图：安全内参 S-01 | Grand View S-06",
    )

    # Slide 5 - Competition
    add_section_slide(
        prs,
        "竞争格局 — 窗口在哪？",
        [
            ("头部厂商：", "奇安信（连续四年第一）、绿盟、深信服、安恒、360"),
            ("专精厂商：", "观成科技 — IDC 2024 加密流量检测代表厂商"),
            ("技术路线：", "AI 不解密（观成/360/绿盟）| 可选解密（奇安信/深信服）| 全量解密"),
            ("我们的机会：", "AI 不解密 + GenAI 降噪 + 云原生东西向流量检测"),
        ],
        footer_note="📸 证据截图：IDC 2024 份额图 S-01/S-08 | 精确 % 需 IDC 报告截图 S-09",
    )

    # Slide 6 - Product direction
    add_section_slide(
        prs,
        "NDR 2.0 怎么做？— 产品能力建议",
        [
            ("必做能力：", "TLS 元数据解析 + JA3/行为特征 + AI 不解密威胁检测"),
            ("增强能力：", "可疑流量选择性解密（平衡隐私与检出率）"),
            ("融合能力：", "GenAI 告警降噪 | API/敏感数据检测 | 云原生全流量"),
            ("竞品对标：", "奇安信、360、深信服、观成科技（PoC 实测待开展）"),
        ],
        footer_note="来源：IDC 2024 技术趋势 + V1.1 竞品公开资料",
    )

    # Slide 7 - Data credibility
    add_table_slide(
        prs,
        "数据可信度 — 回应「数据要有依据」",
        ["类别", "数量", "说明"],
        [
            ["✅ 已公开验证", "18 条", "附来源链接（IDC/Zscaler/HTTP Archive 等）"],
            ["⚠️ 已修正", "10 处", "V1.0 口径/年份/数值错误已更正"],
            ["❌ 已删除", "3 处", "CR5、ETD 32% CAGR、金融渗透率 68% 无出处"],
            ["🔴 待补充", "3 项", "IDC 精确份额截图 | 竞品 PoC | ETD 细分量化"],
        ],
        footer_note="交付物：《数据证据索引与核验清单》+ 报告 V1.1",
    )

    # Slide 8 - Next steps
    add_section_slide(
        prs,
        "决策建议 & 下一步",
        [
            ("请领导决策：", "批准 NDR 2.0 加密流量分析纳入版本规划"),
            ("Q3 行动 ①：", "5 家竞品加密流量检测 PoC（奇安信/360/深信服/Palo Alto/Fortinet）"),
            ("Q3 行动 ②：", "补充 5 张核心数据截图，嵌入汇报材料"),
            ("Q3 行动 ③：", "输出加密流量分析需求 PRD 初稿"),
            ("需支持：", "竞品测试环境；IDC 报告采购（可选，获取精确份额）"),
        ],
        footer_note="第一阶段目标：定性验证通过 → 第二阶段：量化 + 竞品实测",
        highlight_first=True,
    )

    out = "/workspace/NDR2.0领导汇报-加密流量分析可行性验证.pptx"
    prs.save(out)
    print(f"Saved: {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
